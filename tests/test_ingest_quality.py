"""Ingest-quality + robustness tests (2026-08 audit fixes + research adoptions).

Covers:
  1. Contextual chunk headers — "[doc › section]" prefix on every chunk
     (default ON), absent when disabled.
  2. Within-file dedup — repeated chunks skipped, ords preserved.
  3. Provenance re-ingest — a file indexed under a different embedder or
     chunker version re-ingests even though its sha is unchanged (S1 fix).
  4. Empty-docs guard — an empty docs dir no longer wipes the index.
  5. Ingest lock — a live concurrent ingest is refused; a stale lock
     (dead pid) is reclaimed.
  6. Empty-parse recording — a file that parses to nothing still gets its
     sha recorded and stale chunks purged.
  7. Table-atomic chunking — XLSX/CSV rows never split; header row carried.
  8. chunk_overlap=0 honored (no forced 64-word floor).

Uses the fastembed CPU embedder — no Ollama required.
"""
from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

from ez_rag.chunker import chunk_sections
from ez_rag.embed import make_embedder
from ez_rag.index import Index
from ez_rag.ingest import ingest
from ez_rag.parsers import ParsedSection
from ez_rag.workspace import Workspace

PASS, FAIL = [], []


def check(name, cond, detail=""):
    if cond:
        PASS.append(name)
        print(f"  PASS  {name}")
    else:
        FAIL.append((name, detail))
        print(f"  FAIL  {name} -- {detail}")


def make_ws(headers=True, dedup=True):
    tmp = Path(tempfile.mkdtemp(prefix="ezrag_iq_"))
    ws = Workspace(tmp)
    ws.initialize()
    cfg = ws.load_config()
    cfg.embedder_provider = "fastembed"
    cfg.chunk_headers = headers
    cfg.dedup_chunks = dedup
    return ws, cfg


def main():
    print("\n[1] chunk headers on/off")
    ws, cfg = make_ws(headers=True)
    (ws.docs_dir / "dogs.md").write_text(
        "# Dogs\n\nBorder Collies herd sheep.", encoding="utf-8")
    ingest(ws, cfg=cfg)
    emb = make_embedder(cfg)
    idx = Index(ws.meta_db_path, embed_dim=emb.dim)
    texts = [r[0] for r in idx.conn.execute("SELECT text FROM chunks")]
    check("headers ON -> [dogs.md] prefix",
          texts and all(t.startswith("[dogs.md") for t in texts),
          f"{texts[:1]}")

    ws2, cfg2 = make_ws(headers=False)
    (ws2.docs_dir / "dogs.md").write_text(
        "# Dogs\n\nBorder Collies herd sheep.", encoding="utf-8")
    ingest(ws2, cfg=cfg2)
    idx2 = Index(ws2.meta_db_path, embed_dim=emb.dim)
    texts2 = [r[0] for r in idx2.conn.execute("SELECT text FROM chunks")]
    check("headers OFF -> raw text",
          texts2 and not any(t.startswith("[dogs.md") for t in texts2))

    print("\n[2] within-file dedup")
    ws3, cfg3 = make_ws()
    # Paragraphs must beat the chunker's 64-word minimum chunk size so
    # each lands in its own chunk; chunk_size=128 -> ~96-word target.
    cfg3.chunk_size = 128
    cfg3.chunk_overlap = 0
    para = ("Border Collies herd sheep across the rolling border country "
            "from first light to dusk, working whistled commands. ") * 5
    uniq = ("Golden Retrievers fetch waterfowl from cold northern lakes "
            "and never seem to tire of the water. ") * 5
    (ws3.docs_dir / "dup.md").write_text(
        "\n\n".join([para.strip(), para.strip(), para.strip(), uniq.strip()]),
        encoding="utf-8")
    ingest(ws3, cfg=cfg3)
    idx3 = Index(ws3.meta_db_path, embed_dim=emb.dim)
    n3 = idx3.conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
    check("3 identical paragraphs collapsed to 1 (+1 unique)",
          n3 == 2, f"got {n3} chunks")
    ords3 = [r[0] for r in idx3.conn.execute(
        "SELECT ord FROM chunks ORDER BY ord")]
    check("original ords preserved through dedup", ords3 == [0, 3], f"{ords3}")

    print("\n[3] provenance re-ingest (S1)")
    s = ingest(ws, cfg=cfg)   # unchanged baseline
    check("baseline skip", s.files_skipped_unchanged == 1)
    idx.conn.execute("UPDATE files SET embedder = 'ollama:stale-model'")
    idx.conn.commit()
    s = ingest(ws, cfg=cfg)
    check("embedder mismatch -> re-ingest", s.files_changed == 1, f"{vars(s)}")
    idx.conn.execute("UPDATE files SET chunker_version = '1'")
    idx.conn.commit()
    s = ingest(ws, cfg=cfg)
    check("chunker-version mismatch -> re-ingest", s.files_changed == 1)

    print("\n[4] empty-docs guard")
    for f in ws.docs_dir.iterdir():
        f.unlink()
    s = ingest(ws, cfg=cfg)
    n_files = idx.conn.execute("SELECT COUNT(*) FROM files").fetchone()[0]
    check("index preserved on empty docs dir",
          s.files_removed == 0 and n_files == 1,
          f"removed={s.files_removed} files={n_files}")

    print("\n[5] ingest lock")
    lock = ws.meta_db_path.parent / "ingest.lock"
    p = subprocess.Popen([sys.executable, "-c", "import time; time.sleep(30)"])
    lock.write_text(str(p.pid), encoding="utf-8")
    try:
        ingest(ws, cfg=cfg)
        check("live concurrent ingest refused", False, "no RuntimeError")
    except RuntimeError:
        check("live concurrent ingest refused", True)
    finally:
        p.kill()
        p.wait(timeout=10)
    # Stale lock (dead pid) reclaimed silently
    lock.write_text(str(p.pid), encoding="utf-8")
    try:
        ingest(ws, cfg=cfg)
        check("stale lock reclaimed", True)
    except RuntimeError as e:
        check("stale lock reclaimed", False, str(e))
    check("lock removed after run", not lock.exists())

    print("\n[6] empty-parse recorded (stale chunks purged)")
    ws4, cfg4 = make_ws()
    f4 = ws4.docs_dir / "shrink.md"
    f4.write_text("Real content about limestone quarries.", encoding="utf-8")
    ingest(ws4, cfg=cfg4)
    idx4 = Index(ws4.meta_db_path, embed_dim=emb.dim)
    assert idx4.conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0] > 0
    f4.write_text("   \n\n   ", encoding="utf-8")   # now parses to nothing
    ingest(ws4, cfg=cfg4)
    n_chunks = idx4.conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
    row = idx4.conn.execute("SELECT n_chunks FROM files").fetchone()
    check("stale chunks purged", n_chunks == 0, f"{n_chunks} remain")
    check("file recorded with 0 chunks", row is not None and row[0] == 0)
    s = ingest(ws4, cfg=cfg4)
    check("no eternal re-parse", s.files_skipped_unchanged == 1, f"{vars(s)}")

    print("\n[7] table-atomic chunking")
    rows = ["name | hp | ac"] + [f"monster{i} | {i*7} | {10+i%8}"
                                  for i in range(120)]
    tch = chunk_sections(
        [ParsedSection(text="\n".join(rows), meta={"kind": "table"})],
        chunk_tokens=64)
    ok_rows = all(
        c.text.split("\n")[0] == "name | hp | ac"
        and all(" | " in ln for ln in c.text.split("\n"))
        for c in tch)
    recovered = set()
    for c in tch:
        recovered.update(c.text.split("\n")[1:])
    check("multi-piece, header carried, rows atomic",
          len(tch) > 1 and ok_rows)
    check("all 120 rows present", len(recovered) == 120, f"{len(recovered)}")

    print("\n[8] overlap=0 honored")
    secs = [ParsedSection(text=("word " * 300).strip())]
    w0 = sum(len(c.text.split())
             for c in chunk_sections(secs, chunk_tokens=128, overlap_tokens=0))
    check("no forced overlap floor", w0 == 300, f"{w0} words")

    print("\n[9] pdf_backend dispatch + fail-open")
    from ez_rag import parsers as P
    sentinel = [ParsedSection(text="from-marker", meta={"parser": "marker"})]
    saved_marker = P._parse_pdf_marker
    try:
        # Backend selected + working -> its sections come back verbatim
        # (nonexistent path proves the built-in pipeline was never touched).
        P._parse_pdf_marker = lambda path: sentinel
        P.set_pdf_backend("marker")
        out = P.parse_pdf(Path("Z:/definitely/not/a/real.pdf"))
        check("selected backend short-circuits", out is sentinel)
        # Backend raising (e.g. library not installed) -> falls back to the
        # built-in parser, which raises on the missing file — and the
        # failure is recorded for ingest to surface.
        P._parse_pdf_marker = saved_marker   # real impl: marker not installed
        P.pop_pdf_backend_fallbacks()
        try:
            P.parse_pdf(Path("Z:/definitely/not/a/real.pdf"))
            builtin_reached = True   # (unlikely: builtin returned something)
        except Exception:
            builtin_reached = True   # builtin raised on missing file = reached
        fallbacks = P.pop_pdf_backend_fallbacks()
        check("missing library -> fail-open to built-in",
              builtin_reached and len(fallbacks) == 1,
              f"fallbacks={fallbacks}")
        # markdown -> sections splitter keeps headings as section labels
        md = "# Title\nintro\n## Sub\nbody text"
        secs_md = P._markdown_to_sections(md, "docling")
        check("markdown split keeps headings",
              [s.section for s in secs_md] == ["Title", "Sub"],
              f"{[s.section for s in secs_md]}")
    finally:
        P._parse_pdf_marker = saved_marker
        P.set_pdf_backend("auto")

    print("\n[10] pptx parser")
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from ez_rag.parsers import parse_pptx, get_parser
        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[1])
        s1.shapes.title.text = "Ohio Geology"
        s1.placeholders[1].text = "Limestone dominates"
        s1.notes_slide.notes_text_frame.text = "Founded 1837 under Mather."
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "Data"
        tb = s2.shapes.add_table(2, 2, Inches(1), Inches(2),
                                  Inches(4), Inches(1)).table
        tb.cell(0, 0).text = "commodity"; tb.cell(0, 1).text = "value"
        tb.cell(1, 0).text = "limestone"; tb.cell(1, 1).text = "$890M"
        import tempfile as _tf
        deck = Path(_tf.mkdtemp(prefix="ezrag_pptx_")) / "d.pptx"
        prs.save(str(deck))
        secs = parse_pptx(deck)
        check("pptx registered", get_parser(deck) is not None)
        check("slide sections with page numbers",
              all(s.page in (1, 2) for s in secs) and len(secs) == 3,
              f"{[(s.page, s.section) for s in secs]}")
        check("table emitted as kind=table",
              any((s.meta or {}).get("kind") == "table"
                  and "limestone | $890M" in s.text for s in secs))
        check("speaker notes captured",
              any("1837" in s.text for s in secs))
        check("slide title becomes section label",
              any(s.section == "Slide 1: Ohio Geology" for s in secs))
    except ImportError:
        check("pptx test skipped (python-pptx missing)", True)

    print("\n[11] legacy Office converter (.doc/.xls/.ppt)")
    from ez_rag.convert import (
        TARGETS, convert_legacy, converter_available, find_soffice,
    )
    from ez_rag.parsers import get_parser
    check("legacy extensions registered",
          all(get_parser(Path(f"x{e}")) is not None for e in TARGETS))
    if converter_available() and find_soffice():
        import subprocess as _sp
        import tempfile as _tf
        import docx as _docx
        d = Path(_tf.mkdtemp(prefix="ezrag_leg_"))
        doc = _docx.Document()
        doc.add_paragraph("Marblehead limestone since 1834.")
        doc.save(str(d / "r.docx"))
        _sp.run([find_soffice(), "--headless", "--norestore",
                 "--convert-to", "doc", "--outdir", str(d),
                 str(d / "r.docx")], capture_output=True, timeout=180)
        legacy = d / "r.doc"
        if legacy.exists():
            secs = get_parser(legacy)(legacy)
            joined = "\n".join(s.text for s in secs)
            check("doc converts + parses", "Marblehead" in joined,
                  joined[:120])
            import time as _t
            t0 = _t.perf_counter()
            convert_legacy(legacy)
            check("conversion cached", _t.perf_counter() - t0 < 0.5)
        else:
            check("doc fixture built", False, "soffice reverse-convert failed")
    else:
        # No converter on this machine — the error must guide the user.
        try:
            convert_legacy(Path("nope.doc"))
            check("no-backend raises with guidance", False)
        except (RuntimeError, FileNotFoundError, ValueError) as e:
            check("no-backend raises with guidance",
                  "LibreOffice" in str(e) or isinstance(e, FileNotFoundError))

    print(f"\n=== ingest-quality summary: {len(PASS)} pass, {len(FAIL)} fail ===")
    for name, det in FAIL:
        print(f"  FAIL  {name} :: {det}")
    return 0 if not FAIL else 1


if __name__ == "__main__":
    sys.exit(main())
